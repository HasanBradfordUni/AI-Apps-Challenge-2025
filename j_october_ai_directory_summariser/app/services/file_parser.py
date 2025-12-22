import os
import hashlib
import mimetypes
from pathlib import Path
import csv
import io
import datetime

# Remove the autoLogger import section and replace with:
try:
    from ..utils.logger_setup import general_logger
except ImportError:
    # Fallback if logger is not available
    print("Logger import failed for file parser, using DummyLogger")
    class DummyLogger:
        def __init__(self, filename): 
            self.file = open(filename, 'a')
            self.file.write("Logging started...\n"+datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")+"\n")
            self.file.close()
        def addToLogs(self, msg): print(f"[LOG] {msg}")
        def addToErrorLogs(self, msg): print(f"[ERROR] {msg}")
        def addToInputLogs(self, prompt, msg): print(f"[INPUT] {prompt}: {msg}")
    general_logger = DummyLogger

# Optional imports for file parsing
try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

class FileParser:
    def __init__(self):
        # Initialize logger
        log_file_path = os.path.join(os.path.dirname(__file__), '..', 'logs', 'file_parser.txt')
        os.makedirs(os.path.dirname(log_file_path), exist_ok=True)
        self.logger = general_logger(log_file_path)
        
        self.text_extensions = ['.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.xml', '.csv']
        self.supported_for_content = ['.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.xml']
        
        available_parsers = []
        if DOCX_AVAILABLE:
            self.supported_for_content.append('.docx')
            available_parsers.append('DOCX')
        if PDF_AVAILABLE:
            self.supported_for_content.append('.pdf')
            available_parsers.append('PDF')
        if EXCEL_AVAILABLE:
            self.supported_for_content.extend(['.xlsx', '.xls'])
            available_parsers.append('Excel')
        if PPTX_AVAILABLE:
            self.supported_for_content.append('.pptx')
            available_parsers.append('PowerPoint')
        
        self.logger.addToLogs(f"FileParser initialized with parsers: {', '.join(available_parsers) if available_parsers else 'basic text only'}")
        self.logger.addToLogs(f"Supported extensions: {', '.join(self.supported_for_content)}")
    
    def analyze_directory_content(self, directory_path):
        """Analyze content of all files in directory"""
        self.logger.addToLogs(f"Starting content analysis for: {directory_path}")
        
        content_analysis = {
            'total_word_count': 0,
            'total_character_count': 0,
            'file_content_stats': {},
            'content_by_type': {},
            'largest_files_by_content': [],
            'parsing_errors': [],
            'supported_files': 0,
            'unsupported_files': 0
        }
        
        file_count = 0
        
        for root, dirs, files in os.walk(directory_path):
            for file in files:
                file_path = os.path.join(root, file)
                self._analyze_file_content(file_path, content_analysis)
                file_count += 1
                
                # Log progress every 50 files
                if file_count % 50 == 0:
                    self.logger.addToLogs(f"Content analysis progress: {file_count} files processed")
        
        # Calculate averages
        if content_analysis['supported_files'] > 0:
            content_analysis['average_words_per_file'] = (
                content_analysis['total_word_count'] / content_analysis['supported_files']
            )
            content_analysis['average_chars_per_file'] = (
                content_analysis['total_character_count'] / content_analysis['supported_files']
            )
            
            self.logger.addToLogs(f"Content analysis complete: {content_analysis['supported_files']} supported files")
            self.logger.addToLogs(f"Total words: {content_analysis['total_word_count']}")
            self.logger.addToLogs(f"Total characters: {content_analysis['total_character_count']}")
            self.logger.addToLogs(f"Average words per file: {content_analysis['average_words_per_file']:.1f}")
        else:
            self.logger.addToLogs("No supported files found for content analysis")
        
        if content_analysis['parsing_errors']:
            self.logger.addToErrorLogs(f"Encountered {len(content_analysis['parsing_errors'])} parsing errors")
        
        # Sort largest files by content
        content_analysis['largest_files_by_content'] = sorted(
            content_analysis['largest_files_by_content'],
            key=lambda x: x['word_count'],
            reverse=True
        )[:10]
        
        return content_analysis
    
    def _analyze_file_content(self, file_path, content_analysis):
        """Analyze content of individual file"""
        try:
            file_ext = Path(file_path).suffix.lower()
            file_size = os.path.getsize(file_path)
            
            if file_ext in self.supported_for_content:
                content = self._extract_text_content(file_path, file_ext)
                
                if content is not None:
                    word_count = len(content.split()) if content else 0
                    char_count = len(content) if content else 0
                    
                    content_analysis['total_word_count'] += word_count
                    content_analysis['total_character_count'] += char_count
                    content_analysis['supported_files'] += 1
                    
                    # Store file stats
                    content_analysis['file_content_stats'][file_path] = {
                        'word_count': word_count,
                        'character_count': char_count,
                        'file_size': file_size,
                        'file_type': file_ext
                    }
                    
                    # Track content by type
                    if file_ext not in content_analysis['content_by_type']:
                        content_analysis['content_by_type'][file_ext] = {
                            'total_words': 0,
                            'total_chars': 0,
                            'file_count': 0
                        }
                    
                    content_analysis['content_by_type'][file_ext]['total_words'] += word_count
                    content_analysis['content_by_type'][file_ext]['total_chars'] += char_count
                    content_analysis['content_by_type'][file_ext]['file_count'] += 1
                    
                    # Track largest files by content
                    if len(content_analysis['largest_files_by_content']) < 50:
                        content_analysis['largest_files_by_content'].append({
                            'path': file_path,
                            'name': os.path.basename(file_path),
                            'word_count': word_count,
                            'character_count': char_count,
                            'file_type': file_ext
                        })
                    
                    # Log significant content files
                    if word_count > 1000:
                        self.logger.addToLogs(f"Large content file found: {os.path.basename(file_path)} ({word_count} words)")
                        
                else:
                    content_analysis['unsupported_files'] += 1
            else:
                content_analysis['unsupported_files'] += 1
                
        except Exception as e:
            self.logger.addToErrorLogs(f"Error analyzing file content {file_path}: {str(e)}")
            content_analysis['parsing_errors'].append({
                'file_path': file_path,
                'error': str(e)
            })
    
    def _extract_text_content(self, file_path, file_ext):
        """Extract text content from various file types"""
        try:
            if file_ext in ['.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.xml', '.csv']:
                return self._read_text_file(file_path)
            elif file_ext == '.docx' and DOCX_AVAILABLE:
                self.logger.addToLogs(f"Reading DOCX file: {os.path.basename(file_path)}")
                return self._read_docx_file(file_path)
            elif file_ext == '.pdf' and PDF_AVAILABLE:
                self.logger.addToLogs(f"Reading PDF file: {os.path.basename(file_path)}")
                return self._read_pdf_file(file_path)
            elif file_ext in ['.xlsx', '.xls'] and EXCEL_AVAILABLE:
                self.logger.addToLogs(f"Reading Excel file: {os.path.basename(file_path)}")
                return self._read_excel_file(file_path)
            elif file_ext == '.pptx' and PPTX_AVAILABLE:
                self.logger.addToLogs(f"Reading PowerPoint file: {os.path.basename(file_path)}")
                return self._read_pptx_file(file_path)
            
        except Exception as e:
            self.logger.addToErrorLogs(f"Error reading {file_path}: {str(e)}")
            return None
        
        return None
    
    def _read_text_file(self, file_path):
        """Read plain text files"""
        encodings = ['utf-8', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    return file.read()
            except UnicodeDecodeError:
                continue
        
        return None
    
    def _read_docx_file(self, file_path):
        """Read DOCX files"""
        if not DOCX_AVAILABLE:
            return None
        
        try:
            doc = docx.Document(file_path)
            text = []
            for paragraph in doc.paragraphs:
                text.append(paragraph.text)
            return '\n'.join(text)
        except Exception:
            return None
    
    def _read_pdf_file(self, file_path):
        """Read PDF files"""
        if not PDF_AVAILABLE:
            return None
        
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = []
                for page in pdf_reader.pages:
                    text.append(page.extract_text())
                return '\n'.join(text)
        except Exception:
            return None
    
    def _read_excel_file(self, file_path):
        """Read Excel files"""
        if not EXCEL_AVAILABLE:
            return None
        
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) if cell is not None else '' for cell in row]
                    text.append('\t'.join(row_text))
            return '\n'.join(text)
        except Exception:
            return None
    
    def _read_pptx_file(self, file_path):
        """Read PowerPoint files"""
        if not PPTX_AVAILABLE:
            return None
        
        try:
            presentation = Presentation(file_path)
            text = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text.append(shape.text)
            return '\n'.join(text)
        except Exception:
            return None
    
    def calculate_content_hash(self, file_path):
        """Calculate hash of file content for duplicate detection"""
        try:
            hash_md5 = hashlib.md5()
            with open(file_path, "rb") as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    hash_md5.update(chunk)
            return hash_md5.hexdigest()
        except Exception:
            return None
    
    def export_to_csv(self, analysis_data):
        """Export analysis data to CSV format"""
        output = io.StringIO()
        writer = csv.writer(output)
        
        # Write summary information
        writer.writerow(['Directory Analysis Summary'])
        writer.writerow(['Directory Path', analysis_data.get('directory_path', '')])
        writer.writerow(['Total Files', analysis_data.get('analysis_result', {}).get('total_files', 0)])
        writer.writerow(['Total Size', analysis_data.get('analysis_result', {}).get('total_size', 0)])
        writer.writerow([])
        
        # Write file type breakdown
        writer.writerow(['File Type Analysis'])
        writer.writerow(['Extension', 'Count', 'Category'])
        
        file_types = analysis_data.get('analysis_result', {}).get('file_types', {})
        for ext, count in file_types.items():
            writer.writerow([ext, count, self._get_file_category(ext)])
        
        writer.writerow([])
        
        # Write content analysis
        writer.writerow(['Content Analysis'])
        content_analysis = analysis_data.get('content_analysis', {})
        writer.writerow(['Total Word Count', content_analysis.get('total_word_count', 0)])
        writer.writerow(['Total Character Count', content_analysis.get('total_character_count', 0)])
        writer.writerow(['Supported Files', content_analysis.get('supported_files', 0)])
        writer.writerow(['Unsupported Files', content_analysis.get('unsupported_files', 0)])
        
        return output.getvalue()
    
    def _get_file_category(self, extension):
        """Categorize file by extension"""
        categories = {
            'documents': ['.pdf', '.doc', '.docx', '.txt', '.md', '.rtf', '.odt'],
            'spreadsheets': ['.xls', '.xlsx', '.csv', '.ods'],
            'presentations': ['.ppt', '.pptx', '.odp'],
            'images': ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.svg', '.tiff'],
            'videos': ['.mp4', '.avi', '.mkv', '.mov', '.wmv', '.flv'],
            'audio': ['.mp3', '.wav', '.flac', '.aac', '.ogg'],
            'archives': ['.zip', '.rar', '.7z', '.tar', '.gz'],
            'code': ['.py', '.js', '.html', '.css', '.java', '.cpp', '.c', '.php']
        }
        
        for category, extensions in categories.items():
            if extension in extensions:
                return category
        return 'other'